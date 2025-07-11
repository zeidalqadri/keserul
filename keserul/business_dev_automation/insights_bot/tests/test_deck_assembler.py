"""
Tests for the DeckAssembler class.

This module contains unit tests for the DeckAssembler class,
focusing on presentation generation and slide creation.
"""

import os
import tempfile
import unittest
from unittest.mock import MagicMock, patch

import pytest

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from ..utils.deck_assembler import DeckAssembler


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not installed")
class TestDeckAssembler(unittest.TestCase):
    """Tests for the DeckAssembler class."""
    
    def setUp(self):
        """Set up test fixtures."""
        self.output_dir = tempfile.mkdtemp()
        self.assembler = DeckAssembler(output_dir=self.output_dir)
        
        # Sample conversation context
        self.context = {
            "variables": {
                "company_name": "Acme Corp",
                "market_segment": "Software",
                "region": "Asia Pacific",
                "time_period": "2023-2025",
                "deck_title": "Market Analysis",
                "primary_color": "#0066CC",
                "secondary_color": "#FF9900"
            },
            "history": [
                {
                    "role": "system",
                    "content": "System message",
                    "timestamp": 1625097600,
                    "metadata": {}
                },
                {
                    "role": "user",
                    "content": "Analyze the software market in Asia Pacific",
                    "timestamp": 1625097610,
                    "metadata": {}
                },
                {
                    "role": "assistant",
                    "content": "# Executive Summary\n\n1. The software market in Asia Pacific is growing rapidly.\n2. Cloud adoption is accelerating.\n3. AI integration is becoming standard.\n4. Local competitors are gaining market share.\n5. Regulatory compliance is a key challenge.",
                    "timestamp": 1625097620,
                    "metadata": {"phase": "phase_2"}
                },
                {
                    "role": "assistant",
                    "content": "# Market Overview\n\nThe software market in Asia Pacific is projected to grow at 15% CAGR.\n\n## Competitive Landscape\n\n- TechCorp (25% market share)\n- InnoSoft (20% market share)\n- LocalTech (15% market share)\n- Others (40% market share)",
                    "timestamp": 1625097630,
                    "metadata": {"phase": "phase_3"}
                },
                {
                    "role": "assistant",
                    "content": "# Customer Analysis\n\n## Customer Segments\n\n1. Enterprise (40% of market)\n2. SMB (35% of market)\n3. Government (15% of market)\n4. Education (10% of market)",
                    "timestamp": 1625097640,
                    "metadata": {"phase": "phase_4"}
                },
                {
                    "role": "assistant",
                    "content": "# Strategic Options\n\nOption 1: Focus on enterprise customers\nOption 2: Expand to SMB segment\nOption 3: Target government contracts\n\n## Recommended Strategy\n\nExpand to SMB segment while maintaining enterprise focus.",
                    "timestamp": 1625097650,
                    "metadata": {"phase": "phase_5"}
                },
                {
                    "role": "assistant",
                    "content": "# Implementation Plan\n\n1. Develop SMB-focused product variants\n2. Establish channel partnerships\n3. Create targeted marketing campaigns\n4. Implement customer success program\n\n## Expected Outcomes\n\n- 25% revenue growth in year 1\n- 35% market share by year 3\n- 40% improvement in customer retention",
                    "timestamp": 1625097660,
                    "metadata": {"phase": "phase_6"}
                }
            ],
            "charts": [
                {
                    "title": "Market Share Distribution",
                    "type": "Pie",
                    "data_series": [
                        {"name": "TechCorp", "values": 25},
                        {"name": "InnoSoft", "values": 20},
                        {"name": "LocalTech", "values": 15},
                        {"name": "Others", "values": 40}
                    ]
                },
                {
                    "title": "Revenue Growth Projection",
                    "type": "Line",
                    "data_series": [
                        {"name": "Baseline", "values": [100, 110, 120, 130, 140]},
                        {"name": "With Strategy", "values": [100, 125, 156, 195, 244]}
                    ],
                    "x_axis": "Years",
                    "y_axis": "Revenue ($ millions)"
                }
            ]
        }
    
    def test_initialization(self):
        """Test initialization of the assembler."""
        self.assertIsNotNone(self.assembler.output_dir)
        self.assertIsNotNone(self.assembler.layouts)
    
    def test_create_presentation(self):
        """Test creating a presentation."""
        presentation = self.assembler.create_presentation(self.context)
        self.assertIsInstance(presentation, Presentation)
        self.assertGreater(len(presentation.slides), 0)
    
    def test_save_presentation(self):
        """Test saving a presentation."""
        presentation = self.assembler.create_presentation(self.context)
        output_path = self.assembler.save_presentation(presentation, "test_deck.pptx")
        self.assertTrue(os.path.exists(output_path))
        self.assertTrue(output_path.endswith("test_deck.pptx"))
    
    def test_extract_phases_from_history(self):
        """Test extracting phases from conversation history."""
        phases = self.assembler._extract_phases_from_history(self.context["history"])
        self.assertIn("phase_2", phases)
        self.assertIn("phase_3", phases)
        self.assertIn("phase_4", phases)
        self.assertIn("phase_5", phases)
        self.assertIn("phase_6", phases)
    
    def test_extract_key_insights(self):
        """Test extracting key insights."""
        content = "# Key Insights\n\n1. First insight\n2. Second insight\n3. Third insight"
        insights = self.assembler._extract_key_insights(content)
        self.assertEqual(len(insights), 3)
        self.assertEqual(insights[0], "First insight")
        self.assertEqual(insights[1], "Second insight")
        self.assertEqual(insights[2], "Third insight")
    
    def test_extract_competitors(self):
        """Test extracting competitors."""
        content = "# Competitive Landscape\n\n- Competitor A\n- Competitor B\n- Competitor C"
        competitors = self.assembler._extract_competitors(content)
        self.assertEqual(len(competitors), 3)
        self.assertEqual(competitors[0], "Competitor A")
        self.assertEqual(competitors[1], "Competitor B")
        self.assertEqual(competitors[2], "Competitor C")
    
    def test_extract_customer_segments(self):
        """Test extracting customer segments."""
        content = "# Customer Segments\n\n1. Segment A\n2. Segment B\n3. Segment C"
        segments = self.assembler._extract_customer_segments(content)
        self.assertEqual(len(segments), 3)
        self.assertEqual(segments[0], "Segment A")
        self.assertEqual(segments[1], "Segment B")
        self.assertEqual(segments[2], "Segment C")
    
    def test_extract_strategic_options(self):
        """Test extracting strategic options."""
        content = "# Strategic Options\n\nOption 1: First option\nOption 2: Second option\nOption 3: Third option"
        options = self.assembler._extract_strategic_options(content)
        self.assertEqual(len(options), 3)
        self.assertIn("First option", options[0])
        self.assertIn("Second option", options[1])
        self.assertIn("Third option", options[2])
    
    def test_extract_recommendation(self):
        """Test extracting recommendation."""
        content = "# Recommended Strategy\n\nThis is the recommended strategy."
        recommendation = self.assembler._extract_recommendation(content)
        self.assertEqual(recommendation, "This is the recommended strategy.")
    
    def test_extract_implementation_steps(self):
        """Test extracting implementation steps."""
        content = "# Implementation\n\n- Step 1\n- Step 2\n- Step 3"
        steps = self.assembler._extract_implementation_steps(content)
        self.assertEqual(len(steps), 3)
        self.assertEqual(steps[0], "Step 1")
        self.assertEqual(steps[1], "Step 2")
        self.assertEqual(steps[2], "Step 3")
    
    def test_extract_expected_outcomes(self):
        """Test extracting expected outcomes."""
        content = "# Expected Outcomes\n\n1. Outcome A\n2. Outcome B\n3. Outcome C"
        outcomes = self.assembler._extract_expected_outcomes(content)
        self.assertEqual(len(outcomes), 3)
        self.assertEqual(outcomes[0], "Outcome A")
        self.assertEqual(outcomes[1], "Outcome B")
        self.assertEqual(outcomes[2], "Outcome C")
    
    def test_find_chart_by_title_keyword(self):
        """Test finding a chart by title keyword."""
        charts = [
            {"title": "Market Share"},
            {"title": "Revenue Growth"},
            {"title": "Customer Segments"}
        ]
        
        # Test finding a chart
        chart = self.assembler._find_chart_by_title_keyword(charts, ["market", "share"])
        self.assertEqual(chart["title"], "Market Share")
        
        # Test not finding a chart
        chart = self.assembler._find_chart_by_title_keyword(charts, ["nonexistent"])
        self.assertIsNone(chart)


@pytest.mark.integration
@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not installed")
class TestDeckAssemblerIntegration:
    """Integration tests for the DeckAssembler class."""
    
    @pytest.fixture
    def setup(self):
        """Set up test fixtures."""
        # Create a temporary directory for output
        output_dir = tempfile.mkdtemp()
        
        # Create the assembler
        assembler = DeckAssembler(output_dir=output_dir)
        
        # Sample conversation context
        context = {
            "variables": {
                "company_name": "Acme Corp",
                "market_segment": "Software",
                "region": "Asia Pacific",
                "time_period": "2023-2025",
                "deck_title": "Market Analysis"
            },
            "history": [
                {
                    "role": "assistant",
                    "content": "# Executive Summary\n\n1. Key insight 1\n2. Key insight 2",
                    "metadata": {"phase": "phase_2"}
                },
                {
                    "role": "assistant",
                    "content": "# Market Overview\n\nMarket details\n\n## Competitive Landscape\n\n- Competitor A\n- Competitor B",
                    "metadata": {"phase": "phase_3"}
                },
                {
                    "role": "assistant",
                    "content": "# Customer Analysis\n\n## Customer Segments\n\n1. Segment A\n2. Segment B",
                    "metadata": {"phase": "phase_4"}
                },
                {
                    "role": "assistant",
                    "content": "# Strategic Options\n\nOption 1: First option\nOption 2: Second option\n\n## Recommended Strategy\n\nThis is the recommendation.",
                    "metadata": {"phase": "phase_5"}
                },
                {
                    "role": "assistant",
                    "content": "# Implementation Plan\n\n1. Step 1\n2. Step 2\n\n## Expected Outcomes\n\n- Outcome A\n- Outcome B",
                    "metadata": {"phase": "phase_6"}
                }
            ],
            "charts": [
                {
                    "title": "Market Share",
                    "type": "Pie"
                },
                {
                    "title": "Revenue Growth",
                    "type": "Line"
                }
            ]
        }
        
        return {
            "assembler": assembler,
            "context": context,
            "output_dir": output_dir
        }
    
    def test_end_to_end_deck_creation(self, setup):
        """Test end-to-end deck creation."""
        assembler = setup["assembler"]
        context = setup["context"]
        
        # Create the presentation
        presentation = assembler.create_presentation(context)
        
        # Save the presentation
        output_path = assembler.save_presentation(presentation, "test_integration.pptx")
        
        # Verify the file exists
        assert os.path.exists(output_path)
        
        # Verify the presentation has slides
        assert len(presentation.slides) > 0
        
        # Verify specific slides exist (title, executive summary, etc.)
        slide_titles = []
        for slide in presentation.slides:
            if hasattr(slide, "shapes") and hasattr(slide.shapes, "title") and slide.shapes.title:
                slide_titles.append(slide.shapes.title.text)
        
        # Check for expected slide titles
        expected_titles = [
            "Market Analysis",  # Title slide
            "Executive Summary",
            "Key Insights",
            "Market Overview",
            "Competitive Landscape",
            "Customer Analysis",
            "Customer Segments",
            "Strategic Options",
            "Recommended Strategy",
            "Implementation Roadmap",
            "Expected Outcomes"
        ]
        
        # Not all titles need to be present, but at least some should be
        matches = [title for title in slide_titles if title in expected_titles]
        assert len(matches) > 0 