"""
Deck Assembler for McKinsey-Insights Bot

This module provides functionality to generate PowerPoint presentations
from the conversation context using python-pptx.
"""

import os
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

# Optional python-pptx import with graceful degradation for test environments
try:
    from pptx import Presentation  # type: ignore
    from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    from pptx.enum.text import PP_ALIGN  # type: ignore
    from pptx.util import Inches, Pt  # type: ignore
    PPTX_AVAILABLE = True
except ImportError:  # pragma: no cover
    PPTX_AVAILABLE = False

    class Presentation:  # type: ignore
        """Minimal stub when python-pptx is unavailable (CI environments)."""
        def __init__(self, *args, **kwargs):
            # Create mock slides collection
            self.slides = type('Slides', (), {'add_slide': self._add_slide_method})()
            # Create indexable slide layouts list
            self.slide_layouts = [None] * 10  # Support indexes 0-9
        
        def _add_slide_method(self, layout):
            """Mock add_slide method for slides collection."""
            slide = type('Slide', (), {
                'shapes': type('Shapes', (), {'title': type('Title', (), {'text': ''})()})(),
                'placeholders': {1: type('Placeholder', (), {'text': '', 'text_frame': type('TextFrame', (), {'text': '', 'paragraphs': [type('Paragraph', (), {'text': '', 'level': 0})()], 'add_paragraph': lambda: type('Paragraph', (), {'text': '', 'level': 0})()})()})()}
            })()
            return slide
        
        def save(self, path: str) -> None:  # noqa: D401
            pass

    class MSO_SHAPE:  # type: ignore
        pass

    class PP_ALIGN:  # type: ignore
        LEFT = 0
        CENTER = 1
        RIGHT = 2

    def Inches(val):  # type: ignore
        return val

    def Pt(val):  # type: ignore
        return val


class DeckAssembler:
    """
    Assembles PowerPoint presentations from conversation context.
    
    This class handles:
    - Creating slides based on conversation content
    - Inserting charts and diagrams
    - Applying consistent formatting and branding
    - Saving the final presentation
    """
    
    def __init__(
        self,
        template_path: Optional[str] = None,
        output_dir: Optional[str] = None
    ):
        """
        Initialize the deck assembler.
        
        Args:
            template_path: Path to the PowerPoint template
            output_dir: Directory to save the generated presentations
        """
        # Set default template path if not provided
        if template_path is None:
            template_path = os.path.join(
                os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                "templates",
                "mckinsey_template.pptx"
            )
        
        # Check if template exists, if not use a blank presentation
        if os.path.exists(template_path):
            self.template_path = template_path
        else:
            self.template_path = None
        
        # Set output directory
        if output_dir is None:
            output_dir = os.path.join(
                os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                "output"
            )
        
        self.output_dir = Path(output_dir)
        os.makedirs(self.output_dir, exist_ok=True)
        
        # Slide layout indices (these will vary based on the template)
        self.layouts = {
            "title": 0,            # Title slide
            "section": 1,          # Section header
            "title_content": 2,    # Title and content
            "title_two_content": 3,  # Title and two content columns
            "title_chart": 4,      # Title and chart
            "title_bullets": 5,    # Title and bullets
            "blank": 6             # Blank slide
        }
    
    def create_presentation(self, context: Dict[str, Any]) -> Presentation:
        """
        Create a new presentation from the conversation context.
        
        Args:
            context: Conversation context dictionary
            
        Returns:
            PowerPoint presentation object
        """
        # Create a new presentation from template or blank
        if self.template_path and os.path.exists(self.template_path):
            prs = Presentation(self.template_path)
        else:
            prs = Presentation()
        
        # Extract variables from context
        variables = context.get("variables", {})
        company_name = variables.get("company_name", "Company")
        market_segment = variables.get("market_segment", "Market")
        region = variables.get("region", "Region")
        time_period = variables.get("time_period", "Period")
        deck_title = variables.get("deck_title", f"{market_segment} Analysis")
        primary_color = variables.get("primary_color", "#0066CC")
        secondary_color = variables.get("secondary_color", "#FF9900")
        
        # Create title slide
        self._add_title_slide(prs, deck_title, company_name, market_segment, region, time_period)
        
        # Process conversation history to create slides
        history = context.get("history", [])
        charts = context.get("charts", [])
        
        # Extract content from conversation phases
        phases = self._extract_phases_from_history(history)
        
        # Create executive summary slide
        if "phase_2" in phases:
            self._add_executive_summary_slide(prs, phases["phase_2"])
        
        # Create market overview slides
        if "phase_3" in phases:
            self._add_market_overview_slides(prs, phases["phase_3"], charts)
        
        # Create customer analysis slides
        if "phase_4" in phases:
            self._add_customer_analysis_slides(prs, phases["phase_4"], charts)
        
        # Create strategic options slides
        if "phase_5" in phases:
            self._add_strategic_options_slides(prs, phases["phase_5"], charts)
        
        # Create implementation plan slides
        if "phase_6" in phases:
            self._add_implementation_plan_slides(prs, phases["phase_6"], charts)
        
        # Create appendix slides
        self._add_appendix_slides(prs, charts)
        
        return prs
    
    def save_presentation(self, presentation: Presentation, filename: Optional[str] = None) -> str:
        """
        Save the presentation to a file.
        
        Args:
            presentation: PowerPoint presentation object
            filename: Filename for the presentation
            
        Returns:
            Path to the saved file
        """
        if filename is None:
            filename = f"mckinsey_deck_{os.urandom(4).hex()}.pptx"
        
        output_path = self.output_dir / filename
        presentation.save(str(output_path))
        
        return str(output_path)
    
    def _extract_phases_from_history(self, history: List[Dict[str, Any]]) -> Dict[str, str]:
        """
        Extract content from different phases in the conversation history.
        
        Args:
            history: Conversation history
            
        Returns:
            Dictionary mapping phase names to content
        """
        phases = {}
        
        for message in history:
            if message["role"] == "assistant":
                metadata = message.get("metadata", {})
                phase = metadata.get("phase")
                
                if phase and phase.startswith("phase_"):
                    if phase not in phases:
                        phases[phase] = message["content"]
                    else:
                        phases[phase] += "\n\n" + message["content"]
        
        return phases
    
    def _add_title_slide(
        self,
        presentation: Presentation,
        title: str,
        company_name: str,
        market_segment: str,
        region: str,
        time_period: str
    ) -> None:
        """
        Add a title slide to the presentation.
        
        Args:
            presentation: PowerPoint presentation object
            title: Slide deck title
            company_name: Company name
            market_segment: Market segment
            region: Geographic region
            time_period: Time period for analysis
        """
        # Add title slide
        slide_layout = presentation.slide_layouts[self.layouts["title"]]
        slide = presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Set subtitle
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = f"{company_name}\n{market_segment} in {region}\n{time_period}"
    
    def _add_executive_summary_slide(self, presentation: Presentation, content: str) -> None:
        """
        Add an executive summary slide to the presentation.
        
        Args:
            presentation: PowerPoint presentation object
            content: Executive summary content
        """
        # Add section header slide
        section_layout = presentation.slide_layouts[self.layouts["section"]]
        section_slide = presentation.slides.add_slide(section_layout)
        section_slide.shapes.title.text = "Executive Summary"
        
        # Extract key insights
        insights = self._extract_key_insights(content)
        
        # Add insights slide
        slide_layout = presentation.slide_layouts[self.layouts["title_bullets"]]
        slide = presentation.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Key Insights"
        
        # Add bullet points
        text_shape = slide.placeholders[1]
        text_frame = text_shape.text_frame
        
        for i, insight in enumerate(insights):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = insight
            p.level = 0
    
    def _add_market_overview_slides(
        self,
        presentation: Presentation,
        content: str,
        charts: List[Dict[str, Any]]
    ) -> None:
        """
        Add market overview slides to the presentation.
        
        Args:
            presentation: PowerPoint presentation object
            content: Market overview content
            charts: Chart specifications
        """
        # Add section header slide
        section_layout = presentation.slide_layouts[self.layouts["section"]]
        section_slide = presentation.slides.add_slide(section_layout)
        section_slide.shapes.title.text = "Market Overview"
        
        # Extract market size and growth
        market_size_chart = self._find_chart_by_title_keyword(charts, ["market size", "growth"])
        
        # Add market size slide
        if market_size_chart:
            slide_layout = presentation.slide_layouts[self.layouts["title_chart"]]
            slide = presentation.slides.add_slide(slide_layout)
            slide.shapes.title.text = market_size_chart.get("title", "Market Size and Growth")
            
            # Add placeholder chart
            chart_placeholder = slide.placeholders[1]
            self._add_chart_placeholder(chart_placeholder, market_size_chart)
        
        # Extract competitive landscape
        competitors = self._extract_competitors(content)
        
        # Add competitive landscape slide
        slide_layout = presentation.slide_layouts[self.layouts["title_content"]]
        slide = presentation.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Competitive Landscape"
        
        # Add competitor list
        text_shape = slide.placeholders[1]
        text_frame = text_shape.text_frame
        
        for i, competitor in enumerate(competitors):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = competitor
            p.level = 0
    
    def _add_customer_analysis_slides(
        self,
        presentation: Presentation,
        content: str,
        charts: List[Dict[str, Any]]
    ) -> None:
        """
        Add customer analysis slides to the presentation.
        
        Args:
            presentation: PowerPoint presentation object
            content: Customer analysis content
            charts: Chart specifications
        """
        # Add section header slide
        section_layout = presentation.slide_layouts[self.layouts["section"]]
        section_slide = presentation.slides.add_slide(section_layout)
        section_slide.shapes.title.text = "Customer Analysis"
        
        # Extract customer segments
        segments = self._extract_customer_segments(content)
        
        # Add customer segments slide
        slide_layout = presentation.slide_layouts[self.layouts["title_bullets"]]
        slide = presentation.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Customer Segments"
        
        # Add bullet points
        text_shape = slide.placeholders[1]
        text_frame = text_shape.text_frame
        
        for i, segment in enumerate(segments):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = segment
            p.level = 0
    
    def _add_strategic_options_slides(
        self,
        presentation: Presentation,
        content: str,
        charts: List[Dict[str, Any]]
    ) -> None:
        """
        Add strategic options slides to the presentation.
        
        Args:
            presentation: PowerPoint presentation object
            content: Strategic options content
            charts: Chart specifications
        """
        # Add section header slide
        section_layout = presentation.slide_layouts[self.layouts["section"]]
        section_slide = presentation.slides.add_slide(section_layout)
        section_slide.shapes.title.text = "Strategic Options"
        
        # Extract strategic options
        options = self._extract_strategic_options(content)
        
        # Add strategic options slide
        slide_layout = presentation.slide_layouts[self.layouts["title_content"]]
        slide = presentation.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Strategic Options"
        
        # Add options
        text_shape = slide.placeholders[1]
        text_frame = text_shape.text_frame
        
        for i, option in enumerate(options):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = option
            p.level = 0
        
        # Add recommendation slide
        recommendation = self._extract_recommendation(content)
        
        if recommendation:
            slide_layout = presentation.slide_layouts[self.layouts["title_content"]]
            slide = presentation.slides.add_slide(slide_layout)
            slide.shapes.title.text = "Recommended Strategy"
            
            # Add recommendation
            text_shape = slide.placeholders[1]
            text_frame = text_shape.text_frame
            text_frame.text = recommendation
    
    def _add_implementation_plan_slides(
        self,
        presentation: Presentation,
        content: str,
        charts: List[Dict[str, Any]]
    ) -> None:
        """
        Add implementation plan slides to the presentation.
        
        Args:
            presentation: PowerPoint presentation object
            content: Implementation plan content
            charts: Chart specifications
        """
        # Add section header slide
        section_layout = presentation.slide_layouts[self.layouts["section"]]
        section_slide = presentation.slides.add_slide(section_layout)
        section_slide.shapes.title.text = "Implementation Plan"
        
        # Extract implementation steps
        steps = self._extract_implementation_steps(content)
        
        # Add implementation steps slide
        slide_layout = presentation.slide_layouts[self.layouts["title_bullets"]]
        slide = presentation.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Implementation Roadmap"
        
        # Add bullet points
        text_shape = slide.placeholders[1]
        text_frame = text_shape.text_frame
        
        for i, step in enumerate(steps):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = step
            p.level = 0
        
        # Add expected outcomes slide
        outcomes = self._extract_expected_outcomes(content)
        
        if outcomes:
            slide_layout = presentation.slide_layouts[self.layouts["title_bullets"]]
            slide = presentation.slides.add_slide(slide_layout)
            slide.shapes.title.text = "Expected Outcomes"
            
            # Add bullet points
            text_shape = slide.placeholders[1]
            text_frame = text_shape.text_frame
            
            for i, outcome in enumerate(outcomes):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = outcome
                p.level = 0
    
    def _add_appendix_slides(self, presentation: Presentation, charts: List[Dict[str, Any]]) -> None:
        """
        Add appendix slides to the presentation.
        
        Args:
            presentation: PowerPoint presentation object
            charts: Chart specifications
        """
        if not charts:
            return
        
        # Add section header slide
        section_layout = presentation.slide_layouts[self.layouts["section"]]
        section_slide = presentation.slides.add_slide(section_layout)
        section_slide.shapes.title.text = "Appendix"
        
        # Add charts that haven't been used yet
        for chart in charts:
            slide_layout = presentation.slide_layouts[self.layouts["title_chart"]]
            slide = presentation.slides.add_slide(slide_layout)
            slide.shapes.title.text = chart.get("title", "Chart")
            
            # Add placeholder chart
            chart_placeholder = slide.placeholders[1]
            self._add_chart_placeholder(chart_placeholder, chart)
    
    def _add_chart_placeholder(self, placeholder, chart_spec: Dict[str, Any]) -> None:
        """
        Add a chart placeholder with chart specification.
        
        Args:
            placeholder: PowerPoint shape placeholder
            chart_spec: Chart specification
        """
        # In a real implementation, this would create an actual chart
        # For now, we'll just add a text box with the chart specification
        text_frame = placeholder.text_frame
        text_frame.text = f"Chart: {chart_spec.get('title', 'Untitled')}\n"
        p = text_frame.add_paragraph()
        p.text = f"Type: {chart_spec.get('type', 'Bar')}"
        
        if "data_series" in chart_spec:
            p = text_frame.add_paragraph()
            p.text = f"Series: {len(chart_spec['data_series'])} data series"
        
        if "description" in chart_spec:
            p = text_frame.add_paragraph()
            p.text = f"Description: {chart_spec['description']}"
    
    def _extract_key_insights(self, content: str) -> List[str]:
        """
        Extract key insights from the content.
        
        Args:
            content: Content to extract insights from
            
        Returns:
            List of key insights
        """
        insights = []
        
        # Look for numbered lists or bullet points
        matches = re.findall(r"\d+\.\s+(.*?)(?=\d+\.\s+|\Z)", content, re.DOTALL)
        if matches:
            insights = [m.strip() for m in matches]
        
        # If no numbered lists, look for bullet points
        if not insights:
            matches = re.findall(r"- (.*?)(?=- |\Z)", content, re.DOTALL)
            if matches:
                insights = [m.strip() for m in matches]
        
        # If still no insights, split by newlines and take non-empty lines
        if not insights:
            lines = content.split("\n")
            insights = [line.strip() for line in lines if line.strip()]
        
        # Limit to 5 insights
        return insights[:5]
    
    def _extract_competitors(self, content: str) -> List[str]:
        """
        Extract competitors from the content.
        
        Args:
            content: Content to extract competitors from
            
        Returns:
            List of competitors
        """
        competitors = []
        
        # Look for sections about competitors
        competitor_section = re.search(
            r"(?:Competitive Landscape|Key players|Competitors)(.*?)(?=##|\Z)",
            content,
            re.DOTALL
        )
        
        if competitor_section:
            section_text = competitor_section.group(1)
            
            # Look for bullet points or numbered lists
            matches = re.findall(r"- (.*?)(?=- |\Z)", section_text, re.DOTALL)
            if matches:
                competitors = [m.strip() for m in matches]
            
            # If no bullet points, try numbered lists
            if not competitors:
                matches = re.findall(r"\d+\.\s+(.*?)(?=\d+\.\s+|\Z)", section_text, re.DOTALL)
                if matches:
                    competitors = [m.strip() for m in matches]
        
        # If no competitors found, return placeholder
        if not competitors:
            competitors = ["Competitor analysis not available"]
        
        return competitors
    
    def _extract_customer_segments(self, content: str) -> List[str]:
        """
        Extract customer segments from the content.
        
        Args:
            content: Content to extract customer segments from
            
        Returns:
            List of customer segments
        """
        segments = []
        
        # Look for sections about customer segments
        segment_section = re.search(
            r"(?:Customer Segments|Customer Analysis)(.*?)(?=##|\Z)",
            content,
            re.DOTALL
        )
        
        if segment_section:
            section_text = segment_section.group(1)
            
            # Look for bullet points or numbered lists
            matches = re.findall(r"- (.*?)(?=- |\Z)", section_text, re.DOTALL)
            if matches:
                segments = [m.strip() for m in matches]
            
            # If no bullet points, try numbered lists
            if not segments:
                matches = re.findall(r"\d+\.\s+(.*?)(?=\d+\.\s+|\Z)", section_text, re.DOTALL)
                if matches:
                    segments = [m.strip() for m in matches]
        
        # If no segments found, return placeholder
        if not segments:
            segments = ["Customer segment analysis not available"]
        
        return segments
    
    def _extract_strategic_options(self, content: str) -> List[str]:
        """
        Extract strategic options from the content.
        
        Args:
            content: Content to extract strategic options from
            
        Returns:
            List of strategic options
        """
        options = []
        
        # Look for sections about strategic options
        options_section = re.search(
            r"(?:Strategic Options|Options)(.*?)(?=##|\Z)",
            content,
            re.DOTALL
        )
        
        if options_section:
            section_text = options_section.group(1)
            
            # Look for option patterns like "Option 1: ..."
            matches = re.findall(r"Option \d+:?\s+(.*?)(?=Option \d+:?|\Z)", section_text, re.DOTALL)
            if matches:
                options = [m.strip() for m in matches]
            
            # If no options found, look for bullet points
            if not options:
                matches = re.findall(r"- (.*?)(?=- |\Z)", section_text, re.DOTALL)
                if matches:
                    options = [m.strip() for m in matches]
        
        # If no options found, return placeholder
        if not options:
            options = ["Strategic options not available"]
        
        return options
    
    def _extract_recommendation(self, content: str) -> str:
        """
        Extract recommendation from the content.
        
        Args:
            content: Content to extract recommendation from
            
        Returns:
            Recommendation text
        """
        # Look for sections about recommendations
        recommendation_section = re.search(
            r"(?:Recommended Strategy|Recommendation)(.*?)(?=##|\Z)",
            content,
            re.DOTALL
        )
        
        if recommendation_section:
            return recommendation_section.group(1).strip()
        
        return "Recommendation not available"
    
    def _extract_implementation_steps(self, content: str) -> List[str]:
        """
        Extract implementation steps from the content.
        
        Args:
            content: Content to extract implementation steps from
            
        Returns:
            List of implementation steps
        """
        steps = []
        
        # Look for sections about implementation
        implementation_section = re.search(
            r"(?:Implementation|Roadmap)(.*?)(?=##|\Z)",
            content,
            re.DOTALL
        )
        
        if implementation_section:
            section_text = implementation_section.group(1)
            
            # Look for bullet points or numbered lists
            matches = re.findall(r"- (.*?)(?=- |\Z)", section_text, re.DOTALL)
            if matches:
                steps = [m.strip() for m in matches]
            
            # If no bullet points, try numbered lists
            if not steps:
                matches = re.findall(r"\d+\.\s+(.*?)(?=\d+\.\s+|\Z)", section_text, re.DOTALL)
                if matches:
                    steps = [m.strip() for m in matches]
        
        # If no steps found, return placeholder
        if not steps:
            steps = ["Implementation steps not available"]
        
        return steps
    
    def _extract_expected_outcomes(self, content: str) -> List[str]:
        """
        Extract expected outcomes from the content.
        
        Args:
            content: Content to extract expected outcomes from
            
        Returns:
            List of expected outcomes
        """
        outcomes = []
        
        # Look for sections about expected outcomes
        outcomes_section = re.search(
            r"(?:Expected Outcomes|Results|Impact)(.*?)(?=##|\Z)",
            content,
            re.DOTALL
        )
        
        if outcomes_section:
            section_text = outcomes_section.group(1)
            
            # Look for bullet points or numbered lists
            matches = re.findall(r"- (.*?)(?=- |\Z)", section_text, re.DOTALL)
            if matches:
                outcomes = [m.strip() for m in matches]
            
            # If no bullet points, try numbered lists
            if not outcomes:
                matches = re.findall(r"\d+\.\s+(.*?)(?=\d+\.\s+|\Z)", section_text, re.DOTALL)
                if matches:
                    outcomes = [m.strip() for m in matches]
        
        # If no outcomes found, return placeholder
        if not outcomes:
            outcomes = ["Expected outcomes not available"]
        
        return outcomes
    
    def _find_chart_by_title_keyword(self, charts: List[Dict[str, Any]], keywords: List[str]) -> Optional[Dict[str, Any]]:
        """
        Find a chart by keywords in its title.
        
        Args:
            charts: List of chart specifications
            keywords: List of keywords to search for
            
        Returns:
            Chart specification or None if not found
        """
        for chart in charts:
            title = chart.get("title", "").lower()
            if any(keyword.lower() in title for keyword in keywords):
                return chart
        
        return None 